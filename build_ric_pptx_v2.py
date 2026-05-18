"""
Build RIC Leadership Presentation v2 — May 2026
Generates a brand-new PowerPoint from scratch using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from lxml import etree
import os

# ── Brand Colors ──────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(50,  35,  189)
SECONDARY = RGBColor(217, 220, 248)
WHITE     = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64,  64,  64)
LIGHT_GRAY= RGBColor(242, 242, 242)
GREEN     = RGBColor(16,  185, 129)
ORANGE    = RGBColor(245, 158, 11)
RED       = RGBColor(192, 0,   0)

OUTPUT_PATH = "/Users/kyle.hogan/Projects/kyles-arena/RIC-Leadership-Presentation-v2-May2026.pptx"

# ── Presentation Setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]

# ── Helper Functions ──────────────────────────────────────────────────────────

def add_slide(prs):
    """Create blank slide with PRIMARY header bar."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    try:
        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = PRIMARY
        hdr.line.fill.background()
    except Exception as e:
        print(f"  [warn] header bar: {e}")
    return slide


def set_title(slide, title, subtitle):
    """Add title and subtitle textboxes over the header bar."""
    try:
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(22)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
    except Exception as e:
        print(f"  [warn] title: {e}")

    try:
        tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.68), Inches(12.7), Inches(0.35))
        tf2 = tb2.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(12)
        run2.font.color.rgb = SECONDARY
        run2.font.name = "Calibri"
    except Exception as e:
        print(f"  [warn] subtitle: {e}")


def add_footer(slide):
    """Add confidential footer."""
    try:
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.25))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "Rapid Innovation Center | eSimplicity | May 2026 | Confidential"
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(150, 150, 150)
        run.font.name = "Calibri"
    except Exception as e:
        print(f"  [warn] footer: {e}")


def add_centered_text(slide, text, top, font_size, color, bold=False, italic=False, width=13.33, left=0):
    try:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = "Calibri"
    except Exception as e:
        print(f"  [warn] centered text '{text[:30]}': {e}")


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    except Exception as e:
        print(f"  [warn] slide bg: {e}")


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width_pt=None):
    """Add rectangle and return shape."""
    try:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if border_color and border_width_pt:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(border_width_pt)
        elif border_color:
            shape.line.color.rgb = border_color
        else:
            shape.line.fill.background()
        return shape
    except Exception as e:
        print(f"  [warn] add_rect: {e}")
        return None


def add_textbox(slide, text, left, top, width, height, font_size=11, color=DARK_GRAY,
                bold=False, italic=False, align=PP_ALIGN.LEFT, word_wrap=True, name="Calibri"):
    try:
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = name
        return tb
    except Exception as e:
        print(f"  [warn] add_textbox '{text[:30]}': {e}")
        return None


def add_textbox_multiline(slide, lines, left, top, width, height,
                           font_size=11, color=DARK_GRAY, bold=False,
                           align=PP_ALIGN.LEFT, line_spacing=1.2):
    """Add a textbox with multiple lines (list of strings)."""
    try:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size = _Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Calibri"
        return tb
    except Exception as e:
        print(f"  [warn] add_textbox_multiline: {e}")
        return None


def set_cell(tbl, row, col, text, bg_color=None, text_color=DARK_GRAY,
             bold=False, font_size=11, align=PP_ALIGN.LEFT):
    """Set table cell text, font, and background color."""
    cell = tbl.cell(row, col)
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = bold
    run.font.name = "Calibri"
    if bg_color:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for sf in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(sf)
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', str(bg_color))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 1: Title...")
slide1 = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide1, PRIMARY)

add_centered_text(slide1, "Rapid Innovation Center",
                  top=2.2, font_size=40, color=WHITE, bold=True)
add_centered_text(slide1, "Leadership Briefing | May 2026",
                  top=3.1, font_size=22, color=SECONDARY)
add_centered_text(slide1, "Accelerating AI-Enabled Innovation Across Programs and Growth",
                  top=3.7, font_size=15, color=WHITE, italic=True)
add_centered_text(slide1, "eSimplicity | Confidential",
                  top=6.8, font_size=11, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BY THE NUMBERS
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 2: By the Numbers...")
slide2 = add_slide(prs)
set_title(slide2, "RIC — By the Numbers", "Since Launch — February through May 2026")
add_footer(slide2)

boxes = [
    ("35",        "Total Initiatives Launched",  True),
    ("13",        "Initiatives Completed",        False),
    ("21",        "Active Initiatives",           True),
    ("30",        "Requests Submitted",           False),
    ("19",        "Artifacts Delivered",          True),
    ("74",        "Consultation Sessions",        False),
    ("4.7 / 5.0", "Avg CSAT Score",              True),
    ("~3.5 Weeks","Avg Idea-to-Delivery",         False),
    ("320",       "Reusable AI Assets",           True),
]

col_lefts = [0.3, 4.3, 8.3]
row_tops  = [1.2, 2.9, 4.6]
bw, bh    = 3.9, 1.6

for idx, (number, label, is_primary) in enumerate(boxes):
    row = idx // 3
    col = idx % 3
    bx  = col_lefts[col]
    by  = row_tops[row]
    fill = PRIMARY if is_primary else SECONDARY
    try:
        add_rect(slide2, bx, by, bw, bh, fill_color=fill)
    except Exception as e:
        print(f"  [warn] box {idx}: {e}")

    num_color   = WHITE    if is_primary else PRIMARY
    label_color = WHITE    if is_primary else DARK_GRAY

    add_textbox(slide2, number, bx, by + 0.18, bw, 0.6,
                font_size=30, color=num_color, bold=True,
                align=PP_ALIGN.CENTER)
    add_textbox(slide2, label, bx, by + 0.85, bw, 0.55,
                font_size=11, color=label_color,
                align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ENGAGEMENT ACTIVITY (SAFE: two separate charts)
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 3: Engagement Activity...")
slide3 = add_slide(prs)
set_title(slide3, "Engagement Activity — Weekly Trends",
          "Feb 23 – May 15, 2026 | 74 total sessions | 19 artifacts delivered")
add_footer(slide3)

weeks        = ["2/23","3/2","3/9","3/16","3/23","3/30","4/6","4/13","4/20","4/27","5/4","5/11"]
bd_sessions  = [2,1,1,1,3,2,4,3,4,0,2,2]
prog_sessions= [2,4,4,4,7,3,6,3,3,2,7,4]
cumulative   = [4,9,14,19,29,34,44,50,57,59,68,74]

# Chart 1 — Bar (Weekly Sessions by Type)
try:
    cd1 = ChartData()
    cd1.categories = weeks
    cd1.add_series("BD/Growth Sessions", bd_sessions)
    cd1.add_series("Program Sessions",   prog_sessions)
    chart1_placeholder = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(3.2), cd1)
    chart1 = chart1_placeholder.chart
    chart1.has_title = True
    chart1.chart_title.text_frame.text = "Weekly Sessions by Type"

    # Color series
    try:
        chart1.series[0].format.fill.solid()
        chart1.series[0].format.fill.fore_color.rgb = PRIMARY
        chart1.series[1].format.fill.solid()
        chart1.series[1].format.fill.fore_color.rgb = SECONDARY
    except Exception as e:
        print(f"  [warn] chart1 colors: {e}")
except Exception as e:
    print(f"  [warn] chart1: {e}")

# Stat label boxes
try:
    add_rect(slide3, 0.3,  4.3, 2.5, 0.35, fill_color=PRIMARY)
    add_textbox(slide3, "74 Total Sessions", 0.3, 4.32, 2.5, 0.32,
                font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide3, 3.1, 4.3, 4.5, 0.35, fill_color=SECONDARY)
    add_textbox(slide3, "25 BD/Growth  ·  49 Program", 3.1, 4.32, 4.5, 0.32,
                font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
except Exception as e:
    print(f"  [warn] stat boxes: {e}")

# Chart 2 — Line (Cumulative)
try:
    cd2 = ChartData()
    cd2.categories = weeks
    cd2.add_series("Cumulative Total", cumulative)
    chart2_placeholder = slide3.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.3), Inches(4.7), Inches(12.7), Inches(2.5), cd2)
    chart2 = chart2_placeholder.chart
    chart2.has_title = True
    chart2.chart_title.text_frame.text = "Cumulative Sessions"

    try:
        chart2.series[0].format.line.color.rgb = PRIMARY
    except Exception as e:
        print(f"  [warn] chart2 color: {e}")
except Exception as e:
    print(f"  [warn] chart2: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHO WE'VE SERVED
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 4: Programs & Agencies Served...")
slide4 = add_slide(prs)
set_title(slide4, "Programs & Agencies Served",
          "Health · Defense & National Security · Civilian verticals")
add_footer(slide4)

# LEFT header
add_rect(slide4, 0.3, 1.2, 6.2, 0.4, fill_color=PRIMARY)
add_textbox(slide4, "Programs Supported", 0.3, 1.25, 6.2, 0.35,
            font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# LEFT content
left_programs = [
    "● T-MSIS (CMS)",
    "● QDAS (CMS)",
    "● NMSC — National Defense",
    "● NHLBI (NIH)",
    "● LPDS (Navy / Marine Corps)",
    "● IDR — Integrated Data Repository (CMS)",
    "● CQP (CMS)",
    "● SBA OIG",
    "● Health BU / PMR Process",
]
add_textbox_multiline(slide4, left_programs, 0.3, 1.65, 6.2, 4.5,
                      font_size=12, color=DARK_GRAY)

# RIGHT header
add_rect(slide4, 6.8, 1.2, 6.2, 0.4, fill_color=PRIMARY)
add_textbox(slide4, "Agencies & Clients", 6.8, 1.25, 6.2, 0.35,
            font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# RIGHT content
right_agencies = [
    "★ Centers for Medicare & Medicaid Services",
    "★ Defense Information Systems Agency (DISA)",
    "★ Natl Marine & Space Command (NMSC)",
    "★ Natl Heart, Lung & Blood Institute (NHLBI)",
    "★ Small Business Administration OIG",
    "★ State of New Mexico",
    "★ Navy Marine Corps (LPDS)",
    "★ Army (Aberdeen Proving Ground)",
    "★ eSimplicity Growth Team",
]
add_textbox_multiline(slide4, right_agencies, 6.8, 1.65, 6.2, 4.5,
                      font_size=12, color=DARK_GRAY)

# Bottom banner
add_rect(slide4, 0.3, 6.3, 12.7, 0.5, fill_color=SECONDARY)
add_textbox(slide4,
            "9 programs supported  |  8+ agencies engaged  |  Health · Defense · Civilian · Corporate",
            0.3, 6.33, 12.7, 0.45,
            font_size=12, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SPEED TO DELIVERY
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 5: Speed to Delivery...")
slide5 = add_slide(prs)
set_title(slide5, "Speed to Delivery",
          "Average idea-to-prototype cycle time: ~25 days (~3.5 weeks)")
add_footer(slide5)

table_rows = 14  # 1 header + 13 data
table_cols = 4
try:
    tbl_shape = slide5.shapes.add_table(
        table_rows, table_cols,
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
    tbl = tbl_shape.table

    # Column widths
    col_widths = [6.5, 2.0, 1.5, 2.7]
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = Inches(w)

    headers = ["Initiative", "Completed", "Days", "Type"]

    # Header row
    for c, h in enumerate(headers):
        set_cell(tbl, 0, c, h, bg_color=PRIMARY, text_color=WHITE, bold=True, font_size=11, align=PP_ALIGN.CENTER)

    # Data rows
    # (initiative, completed, days, type, days_color)
    data = [
        ("Nextgen Sim - Prima Marketing",   "Mar 27", "<1 day",  "BD/Growth", GREEN),
        ("LPDS AI Tools Consultation",      "Mar 11", "2 days",  "Program",   GREEN),
        ("CMS MESH AI Write-up",            "Apr 3",  "3 days",  "BD/Growth", GREEN),
        ("Aberdeen APBI Demo",              "Apr 10", "8 days",  "BD/Growth", GREEN),
        ("DBX Apps POC - Use Case 4",       "May 8",  "14 days", "BD/Growth", GREEN),
        ("QDAS DBX Apps Research",          "Mar 19", "15 days", "Both",      GREEN),
        ("CMS MESH Salesforce/M365",        "Apr 17", "15 days", "BD/Growth", GREEN),
        ("SBA OIG AI Innovation",           "Apr 24", "30 days", "BD/Growth", ORANGE),
        ("DISA AES RFP AI Assistance",      "Mar 31", "40 days", "BD/Growth", PRIMARY),
        ("DBX Apps POC - Use Case 3",       "May 1",  "43 days", "Both",      PRIMARY),
        ("T-MSIS AI Acceleration",          "Mar 23", "45 days", "Program",   PRIMARY),
        ("New Mexico RFP Demo",             "May 12", "53 days", "BD/Growth", PRIMARY),
        ("NMSC Spectrum Analyst Agent",     "Apr 2",  "55 days", "Program",   PRIMARY),
    ]

    for r, (init, comp, days, typ, days_clr) in enumerate(data):
        row_bg = WHITE if r % 2 == 0 else SECONDARY
        set_cell(tbl, r+1, 0, init,  bg_color=row_bg, font_size=11)
        set_cell(tbl, r+1, 1, comp,  bg_color=row_bg, font_size=11, align=PP_ALIGN.CENTER)
        set_cell(tbl, r+1, 2, days,  bg_color=row_bg, text_color=days_clr, bold=True, font_size=11, align=PP_ALIGN.CENTER)
        set_cell(tbl, r+1, 3, typ,   bg_color=row_bg, font_size=11, align=PP_ALIGN.CENTER)

except Exception as e:
    print(f"  [warn] slide5 table: {e}")

# Note below table
add_textbox(slide5,
            "Shorter = consultations & proposals. Longer = full POC/prototype builds.",
            0.3, 6.85, 12.7, 0.35,
            font_size=10, color=RGBColor(150,150,150), italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROGRAM WINS
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 6: Program Innovation Wins...")
slide6 = add_slide(prs)
set_title(slide6, "Program Innovation Wins",
          "AI delivering measurable impact on active federal programs")
add_footer(slide6)

def make_card(slide, card_left, card_top, card_w, card_h,
              accent_color, title_text, body_lines, quote_text, tag_text=None):
    try:
        # Card background
        add_rect(slide, card_left, card_top, card_w, card_h,
                 fill_color=WHITE, border_color=PRIMARY, border_width_pt=2)
        # Top accent bar
        add_rect(slide, card_left, card_top, card_w, 0.25, fill_color=accent_color)
        # Title
        add_textbox(slide, title_text, card_left + 0.1, card_top + 0.33, card_w - 0.2, 0.42,
                    font_size=13, color=PRIMARY, bold=True)
        # Body
        body_text = "\n".join(body_lines)
        add_textbox(slide, body_text, card_left + 0.1, card_top + 0.78, card_w - 0.2, 2.75,
                    font_size=11, color=DARK_GRAY, word_wrap=True)
        # Quote box
        q_top = card_top + 3.63
        q_h   = 1.15
        add_rect(slide, card_left + 0.1, q_top, card_w - 0.2, q_h, fill_color=SECONDARY)
        add_textbox(slide, quote_text, card_left + 0.15, q_top + 0.07, card_w - 0.3, q_h - 0.1,
                    font_size=10, color=DARK_GRAY, italic=True, word_wrap=True)
        if tag_text:
            add_rect(slide, card_left + 0.1, card_top + card_h - 0.45,
                     card_w - 0.2, 0.38, fill_color=SECONDARY)
            add_textbox(slide, tag_text,
                        card_left + 0.15, card_top + card_h - 0.43,
                        card_w - 0.3, 0.35,
                        font_size=10, color=PRIMARY, bold=True)
    except Exception as e:
        print(f"  [warn] make_card: {e}")


# LEFT CARD
make_card(slide6,
    card_left=0.2, card_top=1.2, card_w=6.4, card_h=5.9,
    accent_color=PRIMARY,
    title_text="T-MSIS — 100%+ Productivity Gain",
    body_lines=[
        "The RIC built an agentic workflow to automate test data generation for",
        "thousands of data validation rules — a process previously done manually.",
        "",
        "● Delivered in days, not weeks",
        "● Adopted and stabilized with the T-MSIS team",
        "● Demonstrated agentic AI on a real federal health program",
        "● Directly supports upcoming T-MSIS recompete",
    ],
    quote_text='"The productivity improvement is at least 100%+! You are truly amazing partners."\n  — Raymond Ling, T-MSIS',
)

# RIGHT CARD
make_card(slide6,
    card_left=6.8, card_top=1.2, card_w=6.3, card_h=5.9,
    accent_color=PRIMARY,
    title_text="NHLBI — Agent Ranker for App Modernization",
    body_lines=[
        "NHLBI needed to prioritize a portfolio of applications for modernization.",
        "The RIC built an Agent Ranker that analyzes profiles and produces",
        "AI-driven recommendations with justification.",
        "",
        "● Replaces manual evaluation process",
        "● Supports ITAC decision-making at NHLBI",
        "● Reusable framework applicable across programs",
        "● Ranking criteria and evaluation metrics in progress",
    ],
    quote_text="",
    tag_text="Program: NHLBI  |  Status: In Progress",
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — BD & PROPOSAL WINS
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 7: BD & Proposal Wins...")
slide7 = add_slide(prs)
set_title(slide7, "BD & Proposal Support Wins",
          "11 growth-supporting initiatives — already exceeding the annual goal of 10+")
add_footer(slide7)

bd_cards = [
    {
        "pos": (0.2, 1.2), "title": "DISA AES — Shaped the Evaluation Criteria",
        "body": ("Built an interactive antenna modeling prototype for the DISA AES proposal. "
                 "DISA then explicitly added innovation as an RFP evaluation factor — "
                 "making our prototype a direct scoring differentiator."),
        "quote": '"Impact increased when customer explicitly included innovation as an eval factor." — Summer Watterson',
    },
    {
        "pos": (6.8, 1.2), "title": "Army Aberdeen — First Army Impression",
        "body": ("Supported a live 2-hour demo at Aberdeen Proving Ground's Meet the Primes event — "
                 "eSimplicity's first Army engagement at this venue. Keith McFarland traveled on-site to deliver."),
        "quote": '"This will play a key role enabling new conversations with the Army." — Don Chipley',
    },
    {
        "pos": (0.2, 4.0), "title": "CMS MESH — Two Gaps Closed Under Pressure",
        "body": ("Two proposal-critical gaps closed: AI narrative written in 3 days, "
                 "Salesforce/M365 integration in 15 days. Both delivered under proposal deadline pressure."),
        "quote": '"Quickly building POCs that allow us to claim experience where gaps exist." — Rajesh Pandian',
    },
    {
        "pos": (6.8, 4.0), "title": "New Mexico — State Government Downselect",
        "body": ("Built a full AI document management demo for NM's water rights system — "
                 "OCR, human-in-the-loop validation, structured data. First state government engagement."),
        "quote": '"It felt like a natural part of the deal lifecycle." — Ian Lowrie',
    },
]

for card in bd_cards:
    cx, cy = card["pos"]
    try:
        cw, ch = 6.3, 2.7
        add_rect(slide7, cx, cy, cw, ch, fill_color=WHITE, border_color=PRIMARY, border_width_pt=1.5)
        add_rect(slide7, cx, cy, cw, 0.2, fill_color=PRIMARY)
        add_textbox(slide7, card["title"], cx + 0.1, cy + 0.25, cw - 0.2, 0.38,
                    font_size=12, color=PRIMARY, bold=True)
        add_textbox(slide7, card["body"], cx + 0.1, cy + 0.65, cw - 0.2, 1.2,
                    font_size=10, color=DARK_GRAY, word_wrap=True)
        # Quote
        q_top = cy + 1.88
        add_rect(slide7, cx + 0.1, q_top, cw - 0.2, 0.72, fill_color=SECONDARY)
        add_textbox(slide7, card["quote"], cx + 0.15, q_top + 0.05, cw - 0.3, 0.62,
                    font_size=9, color=DARK_GRAY, italic=True, word_wrap=True)
    except Exception as e:
        print(f"  [warn] bd card: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — QDAS DATABRICKS TRACK
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 8: QDAS Databricks Track...")
slide8 = add_slide(prs)
set_title(slide8, "QDAS — Sustained Databricks Innovation Track",
          "3 initiatives · 4 POCs · Directly improving QDAS recompete pWIN")
add_footer(slide8)

dbx_boxes = [
    {"left": 0.2, "title": "RIC-102",        "sub": "DBX Research",
     "body": "Evaluated Databricks LLM/AI capabilities. Delivered whitepaper on chatbot feasibility for QDAS users.",
     "tag": "✓ Done"},
    {"left": 3.3, "title": "Use Cases 1 & 2", "sub": "DBX Apps Initial",
     "body": "Proved initial Databricks Apps capability. Evaluated native AI interfaces for non-technical users.",
     "tag": "✓ Done"},
    {"left": 6.4, "title": "RIC-164",         "sub": "Use Case 3 - Portal",
     "body": "Ran CCSQ Portal in Databricks Apps. Validated microservices + Lakebase for persistent storage.",
     "tag": "✓ Done"},
    {"left": 9.5, "title": "RIC-329",         "sub": "Use Case 4 - Deploy",
     "body": "Proved Auto App Deployment for non-technical data consumers. Extends platform to wider audience.",
     "tag": "✓ Done"},
]

for b in dbx_boxes:
    try:
        bw2, bh2, bt = 2.8, 4.2, 1.5
        add_rect(slide8, b["left"], bt, bw2, bh2, fill_color=WHITE, border_color=PRIMARY, border_width_pt=2)
        add_rect(slide8, b["left"], bt, bw2, 0.2, fill_color=PRIMARY)
        add_textbox(slide8, b["title"], b["left"]+0.1, bt+0.28, bw2-0.2, 0.4,
                    font_size=11, color=PRIMARY, bold=True)
        add_textbox(slide8, b["sub"], b["left"]+0.1, bt+0.7, bw2-0.2, 0.35,
                    font_size=10, color=RGBColor(120,120,120))
        add_textbox(slide8, b["body"], b["left"]+0.1, bt+1.1, bw2-0.2, 2.1,
                    font_size=9, color=DARK_GRAY, word_wrap=True)
        add_textbox(slide8, b["tag"], b["left"]+0.1, bt+3.65, bw2-0.2, 0.38,
                    font_size=10, color=GREEN, bold=True)
    except Exception as e:
        print(f"  [warn] dbx box: {e}")

# Arrows
for ax in [3.1, 6.2, 9.3]:
    add_textbox(slide8, "→", ax, 3.2, 0.2, 0.45,
                font_size=20, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# Bottom banner
add_rect(slide8, 0.2, 5.85, 12.9, 0.5, fill_color=PRIMARY)
add_textbox(slide8,
            "Databricks selected as RIC's preferred agentic platform  |  "
            "eSimplicity depth demonstrated across 4 proof points  |  "
            "Directly improving QDAS recompete pWIN",
            0.2, 5.88, 12.9, 0.45,
            font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CSAT
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 9: Customer Satisfaction...")
slide9 = add_slide(prs)
set_title(slide9, "Customer Satisfaction",
          "Average CSAT: 4.7 / 5.0 across 10 scored engagements")
add_footer(slide9)

# Left stat
add_textbox(slide9, "4.7", 0.3, 1.4, 3.5, 2.0,
            font_size=72, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide9, "out of 5.0", 0.3, 3.3, 3.5, 0.5,
            font_size=18, color=PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide9, "10 responses scored", 0.3, 3.8, 3.5, 0.4,
            font_size=12, color=RGBColor(120,120,120), align=PP_ALIGN.CENTER)

# Green callout
add_rect(slide9, 0.3, 4.4, 3.5, 0.6, fill_color=GREEN)
add_textbox(slide9, "7 more responses pending", 0.3, 4.45, 3.5, 0.5,
            font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# CSAT Table
csat_data = [
    ("T-MSIS AI Acceleration",  "Raymond Ling",     "5 ★★★★★", "Program"),
    ("DISA AES RFP",            "Summer Watterson",  "5 ★★★★★", "Growth"),
    ("DBX Apps UC 1&2",         "John Fukuda",       "4 ★★★★",  "Growth"),
    ("NMSC NTIA Redbook",       "Jon Lanenga",       "5 ★★★★★", "Program"),
    ("CMS MESH POC",            "Alisha Hutson",     "5 ★★★★★", "Growth"),
    ("CMS MESH AI Write-Up",    "Rajesh Pandian",    "5 ★★★★★", "Growth"),
    ("Army APBI Demo",          "Don Chipley",       "5 ★★★★★", "Growth"),
    ("DBX Apps UC 3",           "Alisha Hutson",     "4 ★★★★",  "Growth"),
    ("DBX Apps UC 1&2",         "Crissy Paleczka",   "5 ★★★★★", "Program"),
    ("New Mexico RFP",          "Ian Lowrie",        "4 ★★★★",  "Growth"),
]

try:
    t2 = slide9.shapes.add_table(11, 4, Inches(4.1), Inches(1.2), Inches(9.0), Inches(5.5)).table
    c2_widths = [4.5, 2.0, 1.0, 1.5]
    for c, w in enumerate(c2_widths):
        t2.columns[c].width = Inches(w)

    hdrs2 = ["Engagement", "Stakeholder", "Score", "Type"]
    for c, h in enumerate(hdrs2):
        set_cell(t2, 0, c, h, bg_color=PRIMARY, text_color=WHITE, bold=True, font_size=11, align=PP_ALIGN.CENTER)

    for r, (eng, stkh, score, typ) in enumerate(csat_data):
        row_bg = WHITE if r % 2 == 0 else SECONDARY
        set_cell(t2, r+1, 0, eng,   bg_color=row_bg, font_size=11)
        set_cell(t2, r+1, 1, stkh,  bg_color=row_bg, font_size=11)
        set_cell(t2, r+1, 2, score, bg_color=row_bg, font_size=11, align=PP_ALIGN.CENTER)
        set_cell(t2, r+1, 3, typ,   bg_color=row_bg, font_size=11, align=PP_ALIGN.CENTER)
except Exception as e:
    print(f"  [warn] csat table: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THE PLATFORM
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 10: The Platform...")
slide10 = add_slide(prs)
set_title(slide10, "The RIC Platform — Prima Delivery",
          "320 reusable AI assets at delivery.alvisprima.com")
add_footer(slide10)

# Big stat boxes
plat_boxes = [
    (0.3,  PRIMARY,   "108", "Agents",  WHITE,    WHITE),
    (4.65, SECONDARY, "140", "Skills",  PRIMARY,  DARK_GRAY),
    (9.0,  PRIMARY,   "72",  "Plugins", WHITE,    WHITE),
]
for pleft, pfill, pnum, plabel, num_clr, lbl_clr in plat_boxes:
    try:
        pborder = PRIMARY if pfill == SECONDARY else None
        bsz = 1.5 if pfill == SECONDARY else None
        b = slide10.shapes.add_shape(1, Inches(pleft), Inches(1.5), Inches(4.0), Inches(2.0))
        b.fill.solid()
        b.fill.fore_color.rgb = pfill
        if pborder:
            b.line.color.rgb = pborder
            b.line.width = Pt(1.5)
        else:
            b.line.fill.background()
        add_textbox(slide10, pnum, pleft, 1.6, 4.0, 1.0,
                    font_size=36, color=num_clr, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide10, plabel, pleft, 2.7, 4.0, 0.55,
                    font_size=14, color=lbl_clr, align=PP_ALIGN.CENTER)
    except Exception as e:
        print(f"  [warn] plat box: {e}")

# Description
add_textbox(slide10,
    "Prima Delivery is eSimplicity's open AI asset library — production-ready agents, "
    "modular skills, and workflow plugins deployable without starting from scratch.",
    0.3, 3.7, 12.7, 0.8,
    font_size=13, color=DARK_GRAY, word_wrap=True)

# Feature boxes
feat_boxes = [
    (0.3,  "● Compatible with OpenCode, Claude Code, Copilot, Cursor, Windsurf"),
    (4.65, "● Organized by domain: Dev, Security, Data & AI, Infrastructure, Docs"),
    (9.0,  "● Tested, versioned, and ready to deploy across programs"),
]
for fleft, ftext in feat_boxes:
    add_rect(slide10, fleft, 4.6, 4.0, 1.5, fill_color=SECONDARY, border_color=PRIMARY, border_width_pt=1)
    add_textbox(slide10, ftext, fleft + 0.1, 4.7, 3.8, 1.3,
                font_size=11, color=DARK_GRAY, word_wrap=True)

# Bottom banner
add_rect(slide10, 0.0, 6.25, 13.33, 0.45, fill_color=PRIMARY)
add_textbox(slide10,
    "4x the target of 75 assets  |  Available to all programs and BD teams today",
    0.0, 6.27, 13.33, 0.42,
    font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — WHAT'S NEXT
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 11: What's Next...")
slide11 = add_slide(prs)
set_title(slide11, "Looking Ahead — Active Priorities",
          "Key initiatives in flight for H2 2026")
add_footer(slide11)

next_cards = [
    {
        "pos": (0.2, 1.2), "title": "SHINE Platform",
        "body": ("Chat-based AI interface for enterprise knowledge discovery. "
                 "MVP in development — CI/CD established, BDD tests complete. "
                 "6-milestone roadmap targeting Q3 2026 launch."),
        "status": "● Milestone 1 In Progress", "status_color": PRIMARY,
    },
    {
        "pos": (6.8, 1.2), "title": "Apex Agent Deployment",
        "body": ("Packaging and deploying proposal agents to T-MSIS and CIO Growth. "
                 "Agents in QA for deployment. Quick Start Guide, GitHub repo, "
                 "and security config in development."),
        "status": "● Deployment Underway", "status_color": ORANGE,
    },
    {
        "pos": (0.2, 4.0), "title": "IDR Research POC / White Paper",
        "body": ("Early capture effort for CMS IDR — impact score 10. "
                 "White paper and POC demonstrating Snowflake/Databricks integration "
                 "for IDR modernization. In leadership review."),
        "status": "● In Review — Blocker Priority", "status_color": RED,
    },
    {
        "pos": (6.8, 4.0), "title": "Growth Opportunities Dashboard",
        "body": ("Agentic workflow for the Growth team to qualify 20-100 daily BD postings "
                 "automatically. Assigns fit ratings, streamlines opportunity assessment. "
                 "Sponsor: Carrie Yang-Johnson."),
        "status": "● In Progress", "status_color": GREEN,
    },
]

for nc in next_cards:
    cx, cy = nc["pos"]
    try:
        cw, ch = 6.3, 2.65
        add_rect(slide11, cx, cy, cw, ch, fill_color=WHITE, border_color=PRIMARY, border_width_pt=2)
        add_rect(slide11, cx, cy, cw, 0.2, fill_color=PRIMARY)
        add_textbox(slide11, nc["title"], cx+0.1, cy+0.28, cw-0.2, 0.38,
                    font_size=13, color=PRIMARY, bold=True)
        add_textbox(slide11, nc["body"], cx+0.1, cy+0.7, cw-0.2, 1.45,
                    font_size=10, color=DARK_GRAY, word_wrap=True)
        add_textbox(slide11, nc["status"], cx+0.1, cy+2.2, cw-0.2, 0.35,
                    font_size=10, color=nc["status_color"], bold=True)
    except Exception as e:
        print(f"  [warn] next card: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PARTNER PROGRAM UPDATES
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 12: Partner Program Updates...")
slide12 = add_slide(prs)
set_title(slide12, "Partner Program Updates",
          "AWS · Databricks · Salesforce · Snowflake | May 2026")
add_footer(slide12)

partner_cards = [
    {
        "pos": (0.2, 1.2), "color": RGBColor(255,153,0),
        "name": "Amazon Web Services", "tier": "Select → Premier",
        "bullets": [
            "● Active Select partner — working toward Premier",
            "● Requires: 20 accredited, 25 tech certs, 3 competencies, 50 opps",
            "● TPP driving completions: 100% on Associate, Pro & Specialty ✓",
            "● Next: Partner Business Plan & Executive Review",
        ],
    },
    {
        "pos": (6.8, 1.2), "color": RGBColor(255,54,33),
        "name": "Databricks", "tier": "Bronze → Gold",
        "bullets": [
            "● Active Bronze partner — targeting Gold tier",
            "● Requires: Silver baseline + DE/ML certs + Brickbuilder specialization",
            "● TPP: 89% Associate & Pro certs | 23% Sales Badges",
            "● Partnership Day: July 8, 2026 at eSimplicity HQ",
        ],
    },
    {
        "pos": (0.2, 4.0), "color": RGBColor(0,161,224),
        "name": "Salesforce", "tier": "Select → Summit",
        "bullets": [
            "● Active Select partner — long-term Summit target",
            "● Requires: Admin, Platform Dev, Consultant certs + outcomes",
            "● CMS MESH Salesforce/M365 integration delivered (May 2026)",
            "● Cert pursuit and co-sell engagement underway",
        ],
    },
    {
        "pos": (6.8, 4.0), "color": RGBColor(41,181,232),
        "name": "Snowflake", "tier": "Registered → Select",
        "bullets": [
            "● Registered partner — building toward Select tier",
            "● Requires: Partner onboarding, SnowPro Core certs, engagement",
            "● IDR white paper (in review) leverages Snowflake capabilities",
            "● Future TPP program under planning",
        ],
    },
]

for pc in partner_cards:
    cx, cy = pc["pos"]
    pc_clr = pc["color"]
    try:
        cw, ch = 6.3, 2.7
        add_rect(slide12, cx, cy, cw, ch, fill_color=WHITE, border_color=pc_clr, border_width_pt=2)
        add_rect(slide12, cx, cy, cw, 0.2, fill_color=pc_clr)
        add_textbox(slide12, pc["name"], cx+0.1, cy+0.28, cw-0.2, 0.38,
                    font_size=12, color=PRIMARY, bold=True)
        add_textbox(slide12, pc["tier"], cx+0.1, cy+0.65, cw-0.2, 0.3,
                    font_size=10, color=RGBColor(120,120,120))
        btext = "\n".join(pc["bullets"])
        add_textbox(slide12, btext, cx+0.1, cy+0.98, cw-0.2, 1.55,
                    font_size=9, color=DARK_GRAY, word_wrap=True)
    except Exception as e:
        print(f"  [warn] partner card: {e}")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TPP INCENTIVES
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 13: TPP Incentives...")
slide13 = add_slide(prs)
set_title(slide13, "Technical Partner Program (TPP) — Incentives",
          "Time-boxed certification initiative · Directly counts toward partner tier advancement")
add_footer(slide13)

# Left panel
add_rect(slide13, 0.2, 1.2, 6.2, 5.8, fill_color=SECONDARY)
add_textbox(slide13, "What is TPP?", 0.3, 1.35, 6.0, 0.38,
            font_size=14, color=PRIMARY, bold=True)

tpp_body = (
    "The TPP incentivizes eSimplicity employees to earn certifications that advance "
    "our partner tiers — unlocking co-sell opportunities, preferred pricing, and "
    "competitive differentiation.\n\n"
    "✓  Time-boxed with clear milestones\n"
    "✓  Certs count toward next partner tier\n"
    "✓  Started with Databricks and AWS\n"
    "✓  Strong employee response — active pursuit\n"
    "✓  Kyle meeting with partners to align\n"
    "✓  Databricks Day: July 8 @ eSimplicity HQ\n"
    "✓  Future expansion: Snowflake, Salesforce"
)
add_textbox(slide13, tpp_body, 0.3, 1.75, 6.0, 5.0,
            font_size=12, color=DARK_GRAY, word_wrap=True)

# Right panel — cert completion
add_textbox(slide13, "Cert Completion Rates", 6.6, 1.35, 6.2, 0.38,
            font_size=14, color=PRIMARY, bold=True)

# Progress bar helper
def add_progress_bar(slide, label, pct_text, pct_float, top, bar_color, pct_color):
    try:
        add_textbox(slide, label, 6.6, top, 4.8, 0.38, font_size=12, color=DARK_GRAY, bold=True)
        add_textbox(slide, pct_text, 11.5, top - 0.05, 0.7, 0.42,
                    font_size=18, color=pct_color, bold=True, align=PP_ALIGN.RIGHT)
        bg_top = top + 0.42
        add_rect(slide, 6.6, bg_top, 6.2, 0.22, fill_color=RGBColor(220,220,220))
        fill_w = 6.2 * pct_float
        if fill_w > 0:
            add_rect(slide, 6.6, bg_top, fill_w, 0.22, fill_color=bar_color)
    except Exception as e:
        print(f"  [warn] progress bar: {e}")

# AWS bar
add_progress_bar(slide13, "AWS — Associate, Pro & Specialty", "100%", 1.0, 1.8, GREEN, GREEN)
add_textbox(slide13, "✓ COMPLETE", 6.6, 2.5, 3.0, 0.3, font_size=10, color=GREEN, bold=True)

# DBX Certs
add_progress_bar(slide13, "Databricks — Associate & Pro", "89%", 0.89, 2.9, PRIMARY, PRIMARY)
add_textbox(slide13, "● IN PROGRESS", 6.6, 3.6, 3.0, 0.3, font_size=10, color=PRIMARY, bold=True)

# DBX Sales
add_progress_bar(slide13, "Databricks — Sales Badges", "23%", 0.23, 4.0, ORANGE, ORANGE)
add_textbox(slide13, "● EARLY STAGE", 6.6, 4.7, 3.0, 0.3, font_size=10, color=ORANGE, bold=True)

# Partnership Day callout
add_rect(slide13, 6.6, 5.1, 6.2, 1.3, fill_color=PRIMARY)
add_textbox(slide13, "Databricks Partnership Day", 6.65, 5.18, 6.1, 0.42,
            font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide13, "July 8, 2026  ·  eSimplicity HQ", 6.65, 5.58, 6.1, 0.35,
            font_size=11, color=SECONDARY, align=PP_ALIGN.CENTER)
add_textbox(slide13, "Build relationships · Discuss pipelines · Advance tier",
            6.65, 5.93, 6.1, 0.35,
            font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Note
add_textbox(slide13,
    "Future TPP: Snowflake and Salesforce programs under consideration",
    6.6, 6.55, 6.2, 0.3,
    font_size=10, color=RGBColor(150,150,150), italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
print("Building Slide 14: Closing...")
slide14 = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide14, PRIMARY)

add_centered_text(slide14, "The RIC is delivering.", top=2.2, font_size=40, color=WHITE, bold=True)
add_centered_text(slide14, "Fast. Measurable. Mission-Aligned.", top=3.1, font_size=22, color=SECONDARY)
add_centered_text(slide14,
    "35 initiatives  ·  13 completed  ·  4.7 CSAT  ·  3.5 week avg delivery",
    top=3.8, font_size=16, color=WHITE)
add_centered_text(slide14, "Questions?", top=4.8, font_size=18, color=WHITE, italic=True)
add_centered_text(slide14, "eSimplicity Rapid Innovation Center  |  May 2026  |  Confidential",
                  top=6.8, font_size=11, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE & VERIFY
# ═══════════════════════════════════════════════════════════════════════════════
print(f"\nSaving to: {OUTPUT_PATH}")
prs.save(OUTPUT_PATH)

file_size = os.path.getsize(OUTPUT_PATH)
print(f"Saved successfully. File size: {file_size:,} bytes ({file_size/1024:.1f} KB)")

# Verify it loads cleanly
print("Verifying file loads without error...")
from pptx import Presentation as _Prs
test_prs = _Prs(OUTPUT_PATH)
print(f"✓ File loads OK — {len(test_prs.slides)} slides found")
