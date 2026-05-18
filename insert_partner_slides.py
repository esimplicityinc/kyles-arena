"""
Insert two new slides into RIC-Leadership-Presentation-May2026.pptx
after slide 11 (before the closing slide 12).
"""

import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Colors ──────────────────────────────────────────────────────────────────
PRIMARY     = RGBColor(0x32, 0x23, 0xBD)   # #3223BD
SECONDARY   = RGBColor(0xD9, 0xDC, 0xF8)   # #D9DCF8
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x40, 0x40, 0x40)   # #404040
LIGHT_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY    = RGBColor(0x80, 0x80, 0x80)
GREEN       = RGBColor(0x10, 0xB9, 0x81)   # #10B981
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B

AWS_COLOR   = RGBColor(0xFF, 0x99, 0x00)
DB_COLOR    = RGBColor(0xFF, 0x36, 0x21)
SF_COLOR    = RGBColor(0x00, 0xA1, 0xE0)
SNOW_COLOR  = RGBColor(0x29, 0xB5, 0xE8)

FOOTER_TEXT = "RIC PM Reports | eSimplicity | May 13, 2026 | Confidential"


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=None):
    """Add a plain rectangle shape with optional fill/border."""
    from pptx.util import Pt as PtUtil
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()   # clear default line first
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a single-paragraph textbox."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          default_size=11, default_color=DARK_GRAY):
    """
    lines = list of dicts:
      { 'text': str, 'size': int, 'bold': bool, 'color': RGBColor, 'italic': bool, 'space_before': int }
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_def in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = line_def.get('align', PP_ALIGN.LEFT)
        if 'space_before' in line_def:
            p.space_before = Pt(line_def['space_before'])

        run = p.add_run()
        run.text = line_def['text']
        run.font.size = Pt(line_def.get('size', default_size))
        run.font.bold = line_def.get('bold', False)
        run.font.italic = line_def.get('italic', False)
        run.font.color.rgb = line_def.get('color', default_color)
        run.font.name = "Calibri"
    return txBox


def add_slide_header(slide, title, subtitle):
    """Add the standard blue header bar with title and subtitle."""
    # Header bar
    add_rect(slide, 0, 0, 13.33, 1.2, fill_color=PRIMARY)
    # Title
    add_textbox(slide, 0.3, 0.1, 12.0, 0.6,
                title, font_size=24, bold=True, color=WHITE)
    # Subtitle
    add_textbox(slide, 0.3, 0.75, 12.0, 0.35,
                subtitle, font_size=14, bold=False, color=SECONDARY)


def add_footer(slide):
    """Add standard footer text."""
    add_textbox(slide, 0.3, 7.2, 12.7, 0.25,
                FOOTER_TEXT, font_size=10, color=MID_GRAY)


def insert_slide_at(prs, slide, index):
    """Move a newly appended slide to the given index position."""
    xml_slides = prs.slides._sldIdLst
    # The slide was just appended, so it's at the end
    slides_list = list(xml_slides)
    new_elem = slides_list[-1]
    xml_slides.remove(new_elem)
    xml_slides.insert(index, new_elem)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Partner Program Updates
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_partner_updates(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # White background (explicit)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Header
    add_slide_header(slide,
        "Partner Program Updates",
        "AWS · Databricks · Salesforce · Snowflake | May 2026")

    # Footer
    add_footer(slide)

    # ── Card data ────────────────────────────────────────────────────────────
    cards = [
        {
            "color":   AWS_COLOR,
            "left":    0.2,
            "top":     1.3,
            "name":    "Amazon Web Services",
            "tier":    "Select → Premier",
            "bullets": [
                "● Active partner — Select tier achieved",
                "● Goal: Premier tier by end of 2026",
                "● Req: 20 accredited individuals, 25 technical certs,\n   3 competencies, 50 launched opportunities",
                "● Next step: Partner Business Plan & Executive Review",
                "● TPP driving cert completions — 100% complete on\n   Associate, Pro & Specialty",
            ],
        },
        {
            "color":   DB_COLOR,
            "left":    6.93,
            "top":     1.3,
            "name":    "Databricks",
            "tier":    "Bronze → Gold",
            "bullets": [
                "● Active partner — Bronze tier",
                "● Goal: Gold tier by end of 2026",
                "● Req: Silver tier baseline + Data Engineer & ML certs\n   + Brickbuilder specialization",
                "● Partnership Day planned: July 8, 2026 at eSimplicity HQ",
                "● TPP driving cert completions — 89% on\n   Associate & Pro certs",
            ],
        },
        {
            "color":   SF_COLOR,
            "left":    0.2,
            "top":     4.05,
            "name":    "Salesforce",
            "tier":    "Select → Summit",
            "bullets": [
                "● Active partner — Select tier",
                "● Goal: Summit tier (long-term)",
                "● Req: Admin certs, Platform Dev I/II, Consultant certs,\n   verified customer outcomes",
                "● Salesforce/M365 integration delivered for CMS MESH\n   proposal (May 2026)",
                "● Cert pursuit and co-sell engagement underway",
            ],
        },
        {
            "color":   SNOW_COLOR,
            "left":    6.93,
            "top":     4.05,
            "name":    "Snowflake",
            "tier":    "Registered → Select",
            "bullets": [
                "● Registered partner — working toward Select",
                "● Goal: Select tier by end of 2026",
                "● Req: Partner onboarding, SnowPro Core certs,\n   customer engagement",
                "● IDR white paper in progress leverages\n   Snowflake capabilities",
                "● Future TPP program being planned",
            ],
        },
    ]

    CARD_W = 6.0
    CARD_H = 2.6

    for card in cards:
        cl = card["left"]
        ct = card["top"]
        cc = card["color"]

        # Card background rectangle (white, colored border)
        add_rect(slide, cl, ct, CARD_W, CARD_H,
                 fill_color=WHITE, line_color=cc, line_width_pt=2)

        # Accent bar at top of card
        add_rect(slide, cl, ct, CARD_W, 0.25, fill_color=cc)

        # Partner name
        add_textbox(slide, cl + 0.12, ct + 0.3, CARD_W - 0.2, 0.32,
                    card["name"], font_size=13, bold=True, color=DARK_GRAY)

        # Tier line
        add_textbox(slide, cl + 0.12, ct + 0.62, CARD_W - 0.2, 0.25,
                    card["tier"], font_size=11, bold=False,
                    color=RGBColor(0x60, 0x60, 0x60))

        # Bullet lines — build as a single multi-line textbox
        bullet_lines = []
        for b in card["bullets"]:
            bullet_lines.append({
                'text': b,
                'size': 9.5,
                'bold': False,
                'color': DARK_GRAY,
            })
        add_multiline_textbox(slide, cl + 0.1, ct + 0.88,
                              CARD_W - 0.15, CARD_H - 0.95,
                              bullet_lines)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TPP Incentives
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_tpp(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # Header
    add_slide_header(slide,
        "Technical Partner Program (TPP) — Incentives",
        "Time-boxed certification initiative · Counts toward partner tier advancement")

    # Footer
    add_footer(slide)

    # ── LEFT SECTION: Program Overview ───────────────────────────────────────
    L_LEFT  = 0.2
    L_TOP   = 1.3
    L_W     = 6.2
    L_H     = 5.6

    # Light secondary background
    add_rect(slide, L_LEFT, L_TOP, L_W, L_H, fill_color=SECONDARY)

    # "What is TPP?" title
    add_textbox(slide, L_LEFT + 0.18, L_TOP + 0.15, L_W - 0.25, 0.38,
                "What is TPP?",
                font_size=14, bold=True, color=PRIMARY)

    # Body text
    body_lines = [
        {
            'text': "The Technical Partner Program incentivizes eSimplicity employees to earn certifications that directly advance our partner tiers — helping us unlock co-sell opportunities, preferred pricing, and competitive differentiation.",
            'size': 11,
            'bold': False,
            'color': DARK_GRAY,
        },
        {'text': "", 'size': 6, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Time-boxed program with clear milestones", 'size': 11, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Certifications count toward next partner tier", 'size': 11, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Started with Databricks and AWS", 'size': 11, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Positive employee response — people are actively pursuing certs", 'size': 11, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Kyle meeting with partners and getting everyone up to speed", 'size': 11, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Planning Databricks Partnership Day — July 8, 2026 at eSimplicity HQ", 'size': 11, 'bold': False, 'color': DARK_GRAY},
        {'text': "\u2713 Future expansion: Snowflake and other partners being planned", 'size': 11, 'bold': False, 'color': DARK_GRAY},
    ]
    add_multiline_textbox(slide, L_LEFT + 0.18, L_TOP + 0.58, L_W - 0.3, L_H - 0.65,
                          body_lines)

    # ── RIGHT SECTION: Certification Progress ────────────────────────────────
    R_LEFT  = 6.6
    R_TOP   = 1.3
    R_W     = 6.5
    BAR_W   = 6.2
    BAR_H   = 0.25

    # Section title
    add_textbox(slide, R_LEFT, R_TOP + 0.15, R_W, 0.38,
                "Cert Completion Rates",
                font_size=14, bold=True, color=PRIMARY)

    # ── AWS Progress Bar ──────────────────────────────────────────────────────
    aws_label_top  = R_TOP + 0.62
    aws_bar_top    = aws_label_top + 0.42
    aws_tag_top    = aws_bar_top + 0.32

    add_textbox(slide, R_LEFT, aws_label_top, 4.8, 0.35,
                "AWS — Associate, Pro & Specialty",
                font_size=12, bold=True, color=DARK_GRAY)
    add_textbox(slide, R_LEFT + 4.8, aws_label_top, 1.4, 0.35,
                "100%",
                font_size=18, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

    # Progress bar bg + fill (100%)
    add_rect(slide, R_LEFT, aws_bar_top, BAR_W, BAR_H, fill_color=LIGHT_GRAY)
    add_rect(slide, R_LEFT, aws_bar_top, BAR_W * 1.0, BAR_H, fill_color=GREEN)

    add_textbox(slide, R_LEFT, aws_tag_top, 2.0, 0.28,
                "\u2713 COMPLETE",
                font_size=10, bold=True, color=GREEN)

    # ── Databricks Associate & Pro Bar ───────────────────────────────────────
    db1_label_top = aws_tag_top + 0.38
    db1_bar_top   = db1_label_top + 0.42
    db1_tag_top   = db1_bar_top + 0.32

    add_textbox(slide, R_LEFT, db1_label_top, 4.8, 0.35,
                "Databricks — Associate & Pro Certs",
                font_size=12, bold=True, color=DARK_GRAY)
    add_textbox(slide, R_LEFT + 4.8, db1_label_top, 1.4, 0.35,
                "89%",
                font_size=18, bold=True, color=PRIMARY, align=PP_ALIGN.RIGHT)

    add_rect(slide, R_LEFT, db1_bar_top, BAR_W, BAR_H, fill_color=LIGHT_GRAY)
    add_rect(slide, R_LEFT, db1_bar_top, BAR_W * 0.89, BAR_H, fill_color=PRIMARY)

    add_textbox(slide, R_LEFT, db1_tag_top, 2.5, 0.28,
                "\u25cf IN PROGRESS",
                font_size=10, bold=True, color=PRIMARY)

    # ── Databricks Sales Badges Bar ───────────────────────────────────────────
    db2_label_top = db1_tag_top + 0.38
    db2_bar_top   = db2_label_top + 0.42
    db2_tag_top   = db2_bar_top + 0.32

    add_textbox(slide, R_LEFT, db2_label_top, 4.8, 0.35,
                "Databricks — Sales Badges",
                font_size=12, bold=True, color=DARK_GRAY)
    add_textbox(slide, R_LEFT + 4.8, db2_label_top, 1.4, 0.35,
                "23%",
                font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

    add_rect(slide, R_LEFT, db2_bar_top, BAR_W, BAR_H, fill_color=LIGHT_GRAY)
    add_rect(slide, R_LEFT, db2_bar_top, BAR_W * 0.23, BAR_H, fill_color=ORANGE)

    add_textbox(slide, R_LEFT, db2_tag_top, 2.5, 0.28,
                "\u25cf EARLY STAGE",
                font_size=10, bold=True, color=ORANGE)

    # ── Callout box ───────────────────────────────────────────────────────────
    callout_top = db2_tag_top + 0.48
    add_rect(slide, R_LEFT, callout_top, 6.2, 1.4, fill_color=PRIMARY)

    callout_lines = [
        {'text': "Next: Databricks Partnership Day",
         'size': 13, 'bold': True, 'color': WHITE},
        {'text': "July 8, 2026 · eSimplicity HQ",
         'size': 12, 'bold': False, 'color': WHITE, 'space_before': 2},
        {'text': "Build relationships · Discuss pipelines · Accelerate tier advancement",
         'size': 11, 'bold': False, 'color': SECONDARY, 'space_before': 4},
    ]
    add_multiline_textbox(slide, R_LEFT + 0.2, callout_top + 0.15,
                          5.8, 1.1, callout_lines)

    # ── Bottom note ───────────────────────────────────────────────────────────
    note_top = callout_top + 1.5
    add_textbox(slide, R_LEFT, note_top, 6.2, 0.3,
                "Future TPP expansion: Snowflake, Salesforce, and other partners under consideration",
                font_size=10, italic=True, color=MID_GRAY)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    path = '/Users/kyle.hogan/Projects/kyles-arena/RIC-Leadership-Presentation-May2026.pptx'
    prs = Presentation(path)

    original_count = len(prs.slides)
    print(f"Original slide count: {original_count}")

    # Build slide 12 (Partner Updates) — appended last, then repositioned
    build_slide_partner_updates(prs)
    insert_slide_at(prs, None, 11)   # move to index 11 (after slide 11)

    # Build slide 13 (TPP Incentives) — appended last, then repositioned
    build_slide_tpp(prs)
    insert_slide_at(prs, None, 12)   # move to index 12 (after partner updates)

    final_count = len(prs.slides)
    print(f"Final slide count: {final_count}")
    assert final_count == original_count + 2, "Expected 2 new slides!"

    # Verify ordering
    for i, slide in enumerate(prs.slides):
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title_text = shape.text_frame.text.strip()[:60]
                break
        print(f"  Slide {i+1}: {title_text}")

    prs.save(path)
    import os
    size = os.path.getsize(path)
    print(f"\nSaved: {path}")
    print(f"File size: {size:,} bytes ({size/1024:.1f} KB)")
    print("Done.")


if __name__ == "__main__":
    main()
